from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_strategy_ppt():
    # 1. 프레젠테이션 설정 (16:9 와이드스크린)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 여백 및 폰트 설정 상수
    MARGIN_LEFT = Inches(1.0)
    MARGIN_RIGHT = Inches(1.0)
    CONTENT_WIDTH = prs.slide_width - MARGIN_LEFT - MARGIN_RIGHT
    FONT_NAME = "LG Smart UI"  # 시스템에 설치된 폰트 사용

    # 슬라이드 마스터 설정 (배경 등 기본 템플릿)
    blank_slide_layout = prs.slide_layouts[6]

    # --- 도움 함수: 텍스트 박스 추가 (세로 배치 및 줄바꿈 자동 처리) ---
    def add_text_box(slide, text, top, height, font_size=14, bold=False, color=RGBColor(0,0,0), level=0):
        text_box = slide.shapes.add_textbox(MARGIN_LEFT, top, CONTENT_WIDTH, height)
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.margin_bottom = Inches(0.05)
        tf.margin_top = Inches(0.05)

        p = tf.paragraphs[0]
        p.text = text
        p.font.name = FONT_NAME
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.LEFT

        # 계층 구조 들여쓰기 설정 (MECE 구조 시각화)
        if level == 1:
            p.level = 0
            p.space_before = Pt(10)
        elif level == 2:
            p.level = 1
            p.space_before = Pt(0)
        elif level == 3:
            p.level = 2
            p.space_before = Pt(0)

        return text_box.top + height

    # --- 도움 함수: 계층형 텍스트 블록 추가 (■, ○, - 구조 처리) ---
    def add_hierarchical_block(slide, start_top, content_lines, title_height=0.4, body_height=1.8):
        current_top = start_top

        # 메인 섹션 (■)
        if content_lines:
            main_text = content_lines[0]
            # ■ 기호 제거 후 텍스트 처리 (필요시 유지, 여기서는 텍스트 자체에 포함된 것으로 간주)
            current_top = add_text_box(slide, main_text, current_top, title_height, font_size=18, bold=True, color=RGBColor(0, 56, 101), level=1)
            current_top += Inches(0.05) # 섹션 간 간격

        # 서브 섹션 (○, -)
        for line in content_lines[1:]:
            if line.startswith("○ "):
                current_top = add_text_box(slide, line, current_top, Inches(0.4), font_size=14, bold=True, color=RGBColor(50, 50, 50), level=2)
            elif line.startswith("- "):
                current_top = add_text_box(slide, line, current_top, Inches(0.35), font_size=11, bold=False, color=RGBColor(80, 80, 80), level=3)

        return current_top + Inches(0.2) # 블록 간 간격

    # ==========================================
    # PAGE 1: 표지 (Title Slide)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)

    # 대외비 마크
    add_text_box(slide1, "■ 대외비 (CONFIDENTIAL)", Inches(0.5), Inches(0.3), font_size=10, bold=True)

    # 메인 타이틀
    title_box = slide1.shapes.add_textbox(MARGIN_LEFT, Inches(2.0), CONTENT_WIDTH, Inches(1.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PFM 사업총괄 상품 카테고리 구조 다각화 추진 전략 보고"
    p.font.name = FONT_NAME
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

    # 부제
    subtitle_box = slide1.shapes.add_textbox(MARGIN_LEFT, Inches(3.2), CONTENT_WIDTH, Inches(0.6))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "3-Tier 맞춤형 운영모델 및 역량·직종별 Matrix 기반의 구조적 원가 혁신 방안"
    p.font.name = FONT_NAME
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.alignment = PP_ALIGN.CENTER

    # 보고 정보 (하단)
    info_box = slide1.shapes.add_textbox(Inches(0), Inches(6.8), prs.slide_width, Inches(0.4))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2026. 06 | 기획관리실 (PFM 사업총괄)"
    p.font.name = FONT_NAME
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.RIGHT

    # ==========================================
    # PAGE 2: 추진 배경 및 시장 환경 분석
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)

    # 슬라이드 타이틀
    add_text_box(slide2, "I. 추진 배경 및 시장 환경 분석 (Background)", Inches(0.5), Inches(0.5), font_size=24, bold=True, color=RGBColor(0, 56, 101))

    # 콘텐츠 1
    content1 = [
        "■ 노란봉투법 도입에 따른 고용 리스크 관리 체계로의 즉각적인 패러다임 시프트가 시급함.",
        "○ 원청의 사용자 책임 범위 확대에 따른 선제적 대응 체계 구축",
        "- 최근 노동법 개정 추이에 따라 하도급 및 협력사(SP) 근로자 영역까지 원청의 책임 범위가 광범위하게 확대되는 추세임.",
        "- 기존 단순 인력 도급 형태(Headcount 기반 계약)는 상시적인 불법파견 시비 및 공동사용자 책임 리스크를 노출함에 따라 원청의 인사 경영권에 직접적인 위협 요인으로 작용함.",
        "- 서비스의 결과물만을 구매하여 정산하는 Lump-sum(총액 정산형) 계약 및 SLA(서비스 수준 평가) 체제로의 사업 구조 전환이 당면 과제임.",
        "○ 거시적 인구 구조 변화 및 기술 숙련도 양극화에 따른 인력 수급 절벽 직면",
        "- 생산가능인구 감소와 고령화 가속화로 인해 FM 현장 근로자의 평균 연령이 60대를 상회하는 등 상시적 구인난이 심화됨.",
        "- 2030 젊은 층의 전문 기술자격증 취득률은 급감하는 반면, 은퇴 세대의 취득 비중이 증가하는 기술 숙련도 미스매치 현상이 확인됨.",
        "- 고정 인력을 현장에 상주시키는 기존 One-size-fits-all 방식의 운영 효율성이 한계에 도달하여 자산별 투입 자원의 차등화가 요구됨."
    ]
    next_top = add_hierarchical_block(slide2, Inches(1.3), content1, title_height=Inches(0.5), body_height=Inches(2.5))

    # ==========================================
    # PAGE 3: 3-Tier 상품 카테고리별 차별화 운영 전략
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)

    # 슬라이드 타이틀
    add_text_box(slide3, "II. 3-Tier 상품 카테고리별 차별화 운영 전략 (Strategy)", Inches(0.5), Inches(0.5), font_size=24, bold=True, color=RGBColor(0, 56, 101))

    # 콘텐츠 1 (Overview)
    intro_text = "■ 고객 자산 가치와 시장 특성에 따라 Premium, Standard, Lite로 자산을 등급화하고 운영 모델을 차별화함."
    next_top = add_text_box(slide3, intro_text, Inches(1.3), Inches(0.5), font_size=18, bold=True, color=RGBColor(0, 56, 101), level=1)
    next_top += Inches(0.1)

    # 콘텐츠 2 (Premium)
    content_premium = [
        "○ [Premium 자산] 핵심수익시장 대상 맞춤 운영을 통한 고수익 극대화",
        "- 적용 대상 및 시장: 그룹 지주사, 주요 제조 계열사 사옥, IDC(인터넷 데이터 센터) 자산 등 최고 가치 매물 영역임.",
        "- 핵심 추진 과제: 버틀러형 고부가가치 전문 기술 인력의 전담 배치 및 거점 상주 체계 고도화 수행임.",
        "- 운영 프로세스: 2중화 설비 인프라가 무력화되는 공통 요인 고장(Common Cause Failure)을 방지하기 위한 정밀 예방 점검 및 특화 SOP 구축임."
    ]
    next_top = add_hierarchical_block(slide3, next_top, content_premium, title_height=Inches(0.4), body_height=Inches(1.2))

    # 콘텐츠 3 (Standard)
    content_standard = [
        "○ [Standard 자산] 주력시장 대상 기술 기반 운영을 통한 안정적 거래선 유지(Retention)",
        "- 적용 대상 및 시장: 1군 빌딩 및 기존 계약 갱신 대상 주력 경쟁 자산 영역임.",
        "- 핵심 추진 과제: 인건비 절감형 기술 솔루션 도입 및 독자적 FM DX(디지털 전환) 솔루션의 현장 적용 확대임.",
        "- 운영 프로세스: EHP 피크 전력 제어 및 IoT 센서 기반의 예측 정비(PdM) 시스템 구축을 통한 보험성 고정비 삭감임."
    ]
    next_top = add_hierarchical_block(slide3, next_top, content_standard, title_height=Inches(0.4), body_height=Inches(1.2))

    # 콘텐츠 4 (Lite)
    content_lite = [
        "○ [Lite 자산] 레드오션 시장 대상 원가 최적화를 통한 외형 규모 및 시장 점유율 확장",
        "- 적용 대상 및 시장: 2군 신규 빌딩 및 가격 경쟁이 극심한 저단가 수주 공략 자산 영역임.",
        "- 핵심 추진 과제: 서비스 품질 관리 기준의 과감한 이원화 적용 및 초저원가 가격 모델 수립임.",
        "- 운영 프로세스: 중소형 빌딩 대상 연면적별 거점 기동대 순찰 체제 전환 및 공정별 외주 위탁 고도화를 통한 고정비 제거임."
    ]
    next_top = add_hierarchical_block(slide3, next_top, content_lite, title_height=Inches(0.4), body_height=Inches(1.2))

    # ==========================================
    # PAGE 4: 역량·직종 Matrix 기반 실행 과제 및 조직 역할
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)

    # 슬라이드 타이틀
    add_text_box(slide4, "III. 역량·직종 Matrix 기반 실행 과제 및 조직 역할 (Execution)", Inches(0.5), Inches(0.5), font_size=24, bold=True, color=RGBColor(0, 56, 101))

    # 콘텐츠 1
    content_exec1 = [
        "■ 서비스 등급별 목표 달성을 위해 3대 역량 요소를 직종별로 결합하고, 전담 조직의 거버넌스를 확립함.",
        "○ [MECE 역량 Matrix] 전문인력·솔루션·프로세스 관점의 직종별(시설·미화·보안) 핵심 과제 도출",
        "- 전문인력 관점: Premium 시설 직종의 초고숙련 마스터 엔지니어 확보 및 Lite 보안·미화 직종의 가변형 인력 풀 운영 체계 정립임.",
        "- 솔루션 관점: Standard 시설 직종의 에너지 최적화 인프라 연동 및 미화·보안 직종의 디지털 모니터링 시스템 구축임.",
        "- 프로세스 관점: 등급별 SLA 평가 표준화 체제 수립 및 하위 등급(Lite)에서 검증된 효율화 성과를 상위 등급(Standard)으로 전이시키는 승급 경로 설계임."
    ]
    next_top = add_hierarchical_block(slide4, Inches(1.3), content_exec1, title_height=Inches(0.5), body_height=Inches(2.2))

    # 콘텐츠 2
    content_exec2 = [
        "○ [조직별 거버넌스] CEO 직할 및 CDO 전담 조직 체계 구축을 통한 신성장 동력 가속화",
        "- CEO 직할 조직: 고정된 물리 자산 관리 중심에서 공간 이동성(Mobility) 및 임직원 경험(Employee Experience) 영역으로 다각화하기 위한 통근 연계 신사업 모델 구체화 및 현장 검증(Pilot)을 전담함.",
        "- CDO 조직: 현장 시니어 근로자의 저하된 IT 문해력을 고려한 원터치 방식의 '샌디앱' 중심 DX 거버넌스 구현임. 현장 인력은 직관적 일상 점검 및 사진 전송 업무에 집중하고, 중앙 백엔드(Back-end) 관제 플랫폼에서 빅데이터 기반 설비 진단 및 통합 제어를 수행하는 운영 이원화를 확립함."
    ]
    next_top = add_hierarchical_block(slide4, next_top, content_exec2, title_height=Inches(0.4), body_height=Inches(2.0))

    # 파일 저장
    file_name = 'C:/Users/Administrator/Desktop/PFM_전략보고서_기획안.pptx'
    prs.save(file_name)
    print(f"파일이 성공적으로 생성되었습니다: {file_name}")

if __name__ == "__main__":
    create_strategy_ppt()

